# tables.py

from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.table import Table, _Cell
from pptx.shapes.base import BaseShape
from typing import List, Dict, Any


class FormattedTable:
    def __init__(
        self,
        slide,
        num_rows: int,
        num_cols: int,
        left: Inches,
        top: Inches,
        width: Inches,
        height: Inches,
        row_heights: List[float],
        cell_formats: Dict[str, Dict[str, Any]],
    ) -> None:
        self.slide = slide
        self.num_rows = num_rows
        self.num_cols = num_cols
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.row_heights = row_heights
        self.cell_formats = cell_formats
        self.table = self._create_table()

    def _create_table(self) -> Table:
        """
        Creates a new table on the slide with the specified dimensions and row heights.
        """
        table_shape: BaseShape = self.slide.shapes.add_table(
            self.num_rows, self.num_cols, self.left, self.top, self.width, self.height
        )
        table: Table = table_shape.table

        total_height_ratio = sum(self.row_heights)
        for idx, height_ratio in enumerate(self.row_heights):
            row_height = int(self.height * (height_ratio / total_height_ratio))
            table.rows[idx].height = row_height

        return table

    def _format_cell(self, cell: _Cell, text: str, format_key: str) -> None:
        """
        Formats a cell with the specified text and cell format.
        """
        cell_format: Dict[str, Any] = self.cell_formats[format_key]

        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*cell_format["background_color"])

        text_frame = cell.text_frame
        text_frame.text = text

        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER

        run = paragraph.runs[0]
        run.font.size = cell_format["font_size"]
        run.font.bold = cell_format["bold"]
        run.font.color.rgb = RGBColor(*cell_format["font_color"])

        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def populate_table(self, data: List[Dict[str, str]]) -> None:
        """
        Populates the table with the provided data and applies cell formatting.
        """
        for col_idx, item in enumerate(data):
            for row_idx, (key, value) in enumerate(item.items()):
                cell = self.table.cell(row_idx, col_idx)
                self._format_cell(cell, value, key)
