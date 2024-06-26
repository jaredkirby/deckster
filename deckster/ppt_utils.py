# ppt_utils.py
from typing import List
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame, _Run, _Paragraph


def configure_text_run(
    run: _Run, font_size: int, bold: bool, color: RGBColor, font_name: str = "Calibri"
) -> None:
    """
    Configures the properties of a text run.

    Args:
        run: The text run to configure.
        font_size: The font size in points.
        bold: Whether the text should be bold.
        color: The font color as an RGBColor object.
        font_name: The font name (default is "Calibri").
    """
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def add_text_to_shape(
    shape: BaseShape,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = RGBColor(0, 0, 0),
) -> None:
    """
    Adds or updates text in a shape with specified font properties.

    Args:
        shape: The shape to add or update text in.
        text: The text to be added or updated.
        font_size: The font size in points (default is 18).
        bold: Whether the text should be bold (default is False).
        color: The font color as an RGBColor object (default is black).
    """
    text_frame = shape.text_frame
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    # Apply formatting to the first run in the paragraph
    run = paragraph.runs[0]
    configure_text_run(run, font_size, bold, color)


def add_bulleted_list(
    shape: BaseShape,
    items: List[str],
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = RGBColor(0, 0, 0),
) -> None:
    """
    Converts the content of a shape into a bulleted list with specified font properties.

    Args:
        shape: The shape to convert to a bulleted list.
        items: A list of strings representing the bullet points.
        font_size: The font size in points (default is 18).
        bold: Whether the text should be bold (default is False).
        color: The font color as an RGBColor object (default is black).
    """
    text_frame = shape.text_frame
    text_frame.clear()  # Clear existing text to start fresh with bullet points
    for item in items:
        paragraph = text_frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0  # Set bullet level (0 for top-level bullets)
        # Apply formatting to the first run in the paragraph
        run = paragraph.runs[0]
        configure_text_run(run, font_size, bold, color)
