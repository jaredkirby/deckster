# main.py

from pptx import Presentation
from pptx.util import Inches
from data_utils import get_presentation_data
from ppt_utils import add_text_to_shape, add_bulleted_list
from tables import FormattedTable
from error_utils import setup_logging, log_error, log_exception, PowerPointError
from config import (
    PPT_TEMPLATE_PATH,
    DATA_FILE_PATH,
    LOG_FILE_PATH,
    TABLE_LEFT,
    TABLE_TOP,
    TABLE_WIDTH,
    TABLE_HEIGHT,
    TABLE_ROW_HEIGHTS,
    CELL_FORMATS,
)


def generate_presentation():
    try:
        # Set up logging
        setup_logging(LOG_FILE_PATH)

        # Load the PowerPoint template
        prs = Presentation(PPT_TEMPLATE_PATH)

        # Get the presentation data from the CSV file
        presentation_data = get_presentation_data(DATA_FILE_PATH)

        # Add content to the slides
        slide1 = prs.slides[0]
        slide2 = prs.slides[1]

        # Slide 1: Add the tactics table
        tactics_table = FormattedTable(
            slide=slide1,
            num_rows=len(CELL_FORMATS),
            num_cols=len(presentation_data["Tactics"]),
            left=TABLE_LEFT,
            top=TABLE_TOP,
            width=TABLE_WIDTH,
            height=TABLE_HEIGHT,
            row_heights=TABLE_ROW_HEIGHTS,
            cell_formats=CELL_FORMATS,
        )
        tactics_table.populate_table(presentation_data["Tactics"])

        # Slide 2: Add the page title, recommendations, and results/learnings
        for shape in slide2.shapes:
            if not shape.has_text_frame:
                continue
            if shape.name == "Page Title":
                add_text_to_shape(shape, presentation_data["Page Title"], font_size=32)
            elif shape.name in ["Recommendation", "Results and Learnings"]:
                add_bulleted_list(shape, presentation_data[shape.name], font_size=16)

        # Save the updated presentation
        prs.save("updated_presentation.pptx")
        print("Presentation generated successfully!")

    except FileNotFoundError as e:
        error_message = f"File not found: {e.filename}"
        log_error(error_message)
        print(error_message)

    except PowerPointError as e:
        log_exception(e)
        print(f"PowerPoint error: {str(e)}")

    except Exception as e:
        log_exception(e)
        print(f"An error occurred: {str(e)}")


if __name__ == "__main__":
    generate_presentation()
